import os
import sys
import uuid
import tempfile
from flask import Flask, render_template, request, send_file, jsonify
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pptx import Presentation
import io

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10MB max

# Vérification de l'installation de Tesseract
try:
    pytesseract.get_tesseract_version()
except EnvironmentError:
    print("Tesseract n'est pas installé ou n'est pas dans le PATH.")
    sys.exit("Erreur: Tesseract requis - installez-le avec 'sudo pacman -S tesseract tesseract-data-fra'")

def extract_text_from_image(image_path):
    try:
        img = Image.open(image_path)
        return pytesseract.image_to_string(img, lang='fra')
    except Exception as e:
        return f"Erreur OCR: {str(e)}"

def extract_text_from_pdf(pdf_path):
    try:
        text = ""
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        return f"Erreur PDF: {str(e)}"

def extract_text_from_ppt(ppt_path):
    try:
        text = ""
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Erreur PPT: {str(e)}"

def process_file(file_path, file_type):
    if file_type == 'image':
        return extract_text_from_image(file_path)
    elif file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type in ['ppt', 'pptx']:
        return extract_text_from_ppt(file_path)
    else:
        return "Format de fichier non supporté"

def get_file_type(filename):
    ext = filename.split('.')[-1].lower()
    if ext in ['png', 'jpg', 'jpeg', 'bmp', 'gif', 'tiff']:
        return 'image'
    elif ext == 'pdf':
        return 'pdf'
    elif ext in ['ppt', 'pptx']:
        return ext
    return None

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'Aucun fichier dans la requête'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'Aucun fichier sélectionné'}), 400
    
    if file:
        # Vérifier le type de fichier
        file_type = get_file_type(file.filename)
        if not file_type:
            return jsonify({'error': 'Format de fichier non supporté'}), 400
        
        # Créer un dossier unique par session
        session_id = str(uuid.uuid4())
        session_folder = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        os.makedirs(session_folder, exist_ok=True)
        
        # Sauvegarder le fichier
        file_path = os.path.join(session_folder, file.filename)
        file.save(file_path)
        
        # Extraire le texte
        extracted_text = process_file(file_path, file_type)
        
        # Sauvegarder le texte
        txt_filename = f"{os.path.splitext(file.filename)[0]}.txt"
        txt_path = os.path.join(session_folder, txt_filename)
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
        
        # Retourner le chemin pour téléchargement
        return jsonify({
            'download_url': f'/download/{session_id}/{txt_filename}',
            'filename': txt_filename
        })
    
    return jsonify({'error': 'Erreur inconnue'}), 500

@app.route('/download/<session_id>/<filename>')
def download(session_id, filename):
    session_folder = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
    txt_path = os.path.join(session_folder, filename)
    
    return send_file(
        txt_path,
        as_attachment=True,
        download_name=filename,
        mimetype='text/plain'
    )

if __name__ == '__main__':
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    app.run(debug=True)